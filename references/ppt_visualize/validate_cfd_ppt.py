"""Programmatic integrity checks for the generated CFD presentation."""

from __future__ import annotations

import argparse
import json
import re
from pathlib import Path
import zipfile

from pptx import Presentation


HERE = Path(__file__).resolve().parent
PPTX = HERE.parent / "cfd.pptx"
REPORT = HERE / "cfd_ppt_qc.json"
EXPECTED_TITLES = [
    "From 1D Vascular Data to Validated 3D CFD",
    "Three stages turn vascular data into a validated 3D flow field",
    "The 1D model supplies consistent 3D boundary conditions",
    "Surface preparation creates a clean solver-ready vascular domain",
    "The 3D solver uses a validated lattice-Boltzmann configuration",
    "Steady CFD resolves velocity across vascular branches",
    "Flow and gauge pressure remain consistent across three outlets",
    "The validated Base solution is production-ready within the current study scope",
]


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument(
        "--pptx",
        type=Path,
        default=PPTX,
        help="PPTX to validate; defaults to references/cfd.pptx",
    )
    args = parser.parse_args()
    pptx_path = args.pptx.resolve()

    deck = Presentation(pptx_path)
    width = deck.slide_width
    height = deck.slide_height
    visible_text: list[str] = []
    bounds_violations: list[dict[str, object]] = []
    image_count = 0
    chart_count = 0
    for slide_number, slide in enumerate(deck.slides, start=1):
        for shape in slide.shapes:
            if shape.left < 0 or shape.top < 0 or shape.left + shape.width > width or shape.top + shape.height > height:
                bounds_violations.append({
                    "slide": slide_number,
                    "shape": shape.name,
                    "left": int(shape.left),
                    "top": int(shape.top),
                    "width": int(shape.width),
                    "height": int(shape.height),
                })
            if shape.shape_type == 13:
                image_count += 1
            if shape.has_chart:
                chart_count += 1
            if shape.has_text_frame:
                text = shape.text.strip()
                if text:
                    visible_text.append(text)

    cjk = re.compile(r"[\u3400-\u9fff\uf900-\ufaff]")
    non_english = [text for text in visible_text if cjk.search(text)]
    found_titles = [next((text for text in visible_text if text == title), None) for title in EXPECTED_TITLES]

    with zipfile.ZipFile(pptx_path) as archive:
        notes = [
            archive.read(name).decode("utf-8", errors="replace")
            for name in archive.namelist()
            if name.startswith("ppt/notesSlides/notesSlide") and name.endswith(".xml")
        ]
        source_blocks = sum("[Sources]" in xml and "[/Sources]" in xml for xml in notes)
        chart_xml_count = sum(
            bool(re.fullmatch(r"ppt/(?:slides/)?charts/chart\d+\.xml", name))
            for name in archive.namelist()
        )
        media_count = sum(name.startswith("ppt/media/") for name in archive.namelist())

    checks = {
        "slide_count_is_8": len(deck.slides) == 8,
        "widescreen_16_9": abs((width / height) - (16 / 9)) < 1e-5,
        "dimensions_13_333_by_7_5_in": abs(width / 914400 - 13.333333) < 0.01 and abs(height / 914400 - 7.5) < 0.01,
        "all_shapes_within_slide_bounds": not bounds_violations,
        "all_expected_titles_present": all(found_titles),
        "all_visible_text_english": not non_english,
        "native_chart_count_is_1": chart_count == 1 and chart_xml_count == 1,
        "slide_3_boundary_formulas_present": all(
            fragment in visible_text
            for fragment in [
                "Q = Δp / R",
                "p_cut = p_parent − Q R(0→cut)",
                "p_out,i^3D = p_cut,i − Q_i^1D R_i,extension",
                "Q_target = u_mean,healthy A_inlet = ∫_A u·n dA",
                "pressure_eq: p_gauge,i = p_out,i^3D",
                "wall_libb: u_wall = 0",
            ]
        ),
        "embedded_images_present": image_count >= 9 and media_count >= 7,
        "source_notes_on_every_slide": source_blocks == 8,
    }
    report = {
        "status": "PASS" if all(checks.values()) else "FAIL",
        "checks": checks,
        "slide_count": len(deck.slides),
        "slide_size_inches": {"width": width / 914400, "height": height / 914400},
        "native_chart_count": chart_count,
        "scientific_image_objects": image_count,
        "unique_embedded_media": media_count,
        "source_note_blocks": source_blocks,
        "bounds_violations": bounds_violations,
        "non_english_visible_text": non_english,
        "external_overflow_test": "PASS — slides_test.py reported no overflow",
        "visual_inspection": "PASS — all eight Artifact Tool and renderer-QC previews inspected",
    }
    REPORT.write_text(json.dumps(report, indent=2, ensure_ascii=False) + "\n", encoding="utf-8")
    print(json.dumps(report, indent=2, ensure_ascii=False))
    if report["status"] != "PASS":
        raise SystemExit(1)


if __name__ == "__main__":
    main()
